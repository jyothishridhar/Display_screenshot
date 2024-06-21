import streamlit as st
from selenium import webdriver
from PIL import Image
from io import BytesIO
import cv2
import numpy as np
from selenium.webdriver.common.by import By
from pptx import Presentation
from pptx.util import Inches
import time
import os

X_OFFSET = 280
Y_OFFSET = -172

def replace_image_in_cropped_area(screenshot, reference_image_path, img_location):
    reference_image = Image.open(reference_image_path)
    reference_width, reference_height = reference_image.size
    reference_image = reference_image.resize((img_location[2], img_location[3]))
    paste_location = (img_location[0] + X_OFFSET, img_location[1] + Y_OFFSET)
    screenshot.paste(reference_image, paste_location)
    modified_screenshot_path = "modified_screenshot.png"
    screenshot.save(modified_screenshot_path)
    return modified_screenshot_path

def find_and_replace_reference_image(driver, reference_image_path):
    reference_image = cv2.imread(reference_image_path, cv2.IMREAD_UNCHANGED)
    reference_height, reference_width, _ = reference_image.shape
    images = driver.find_elements(By.TAG_NAME, 'iframe')
    screenshot = driver.get_screenshot_as_png()
    screenshot = Image.open(BytesIO(screenshot))
    screenshot = np.array(screenshot)

    for img in images:
        img_width = img.size['width']
        img_height = img.size['height']
        if img_width == reference_width and img_height == reference_height:
            driver.execute_script("arguments[0].scrollIntoView({behavior: 'auto', block: 'center'});", img)
            time.sleep(1)
            screenshot = driver.get_screenshot_as_png()
            screenshot = Image.open(BytesIO(screenshot))
            screenshot = np.array(screenshot)
            screenshot_gray = cv2.cvtColor(screenshot, cv2.COLOR_BGR2GRAY)
            reference_gray = cv2.cvtColor(reference_image, cv2.COLOR_BGR2GRAY)
            result = cv2.matchTemplate(screenshot_gray, reference_gray, cv2.TM_CCOEFF_NORMED)
            _, _, _, max_loc = cv2.minMaxLoc(result)
            modified_screenshot_path = replace_image_in_cropped_area(
                Image.fromarray(screenshot),
                reference_image_path,
                (*max_loc, reference_width, reference_height)
            )
            return modified_screenshot_path
    return None

def save_screenshot_to_pptx(screenshot_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(-0.8)
    top = Inches(0)
    img = Image.open(screenshot_path)
    width, height = img.size
    img_path = "static_screenshot.png"
    img.save(img_path)
    slide.shapes.add_picture(img_path, left, top, width=Inches(width / 80), height=Inches(height / 80))
    pptx_path = "modified_screenshot.pptx"
    prs.save(pptx_path)
    return pptx_path

def main():
    st.title("Image Replacement and PPTX Generator")

    url = st.text_input("Enter the URL:")
    reference_image_file = st.file_uploader("Upload the reference image", type=["png", "jpg", "jpeg"])

    if st.button("Process"):
        if url and reference_image_file:
            reference_image_path = "reference_image.png"
            with open(reference_image_path, "wb") as f:
                f.write(reference_image_file.read())

            options = webdriver.ChromeOptions()
            options.add_argument("--headless")
            options.add_argument("--disable-gpu")
            options.add_argument("--no-sandbox")
            driver = webdriver.Chrome(options=options)
            driver.get(url)
            driver.maximize_window()
            time.sleep(15)

            modified_screenshot_path = find_and_replace_reference_image(driver, reference_image_path)

            if modified_screenshot_path:
                # Display modified screenshot
                st.image(Image.open(modified_screenshot_path), caption="Modified Screenshot", use_column_width=True)
                
                # Save screenshot to PPTX and display download button
                pptx_path = save_screenshot_to_pptx(modified_screenshot_path)
                with open(pptx_path, "rb") as f:
                    st.download_button("Download PowerPoint", f, file_name="modified_screenshot.pptx")

                st.write("PPTX saved at:", pptx_path)  # Debugging line
            else:
                st.write("Reference image not found in the screenshot.")
            driver.quit()

if __name__ == "__main__":
    main()
